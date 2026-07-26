import base64
import hashlib
import json
import logging
import os
import re
import struct
import subprocess
import threading
import time
from datetime import datetime, timezone
from pathlib import Path
from zipfile import ZipFile, BadZipFile

from vant.utils import run_hidden

import requests


DEFAULT_DLP_EXTENSIONS = {
    ".txt", ".log", ".csv", ".json", ".xml", ".md",
    ".doc", ".docx", ".docm", ".rtf",
    ".xls", ".xlsx", ".xlsm",
    ".ppt", ".pptx", ".pptm",
    ".pdf",
    ".odt", ".ods", ".odp",
    ".ini", ".conf", ".cfg", ".yaml", ".yml",
    ".ps1", ".bat", ".cmd", ".sql",
    ".env", ".properties",
    ".html", ".htm", ".eml", ".msg",
}

DEFAULT_DLP_KEYWORDS = [
    {"name": "Clasificado", "classification": "clasificado", "severity": "critical",
     "match_type": "keyword", "pattern": "informacion clasificada",
     "tags": ["estado", "clasificado"]},
    {"name": "Confidential", "classification": "confidential", "severity": "critical",
     "match_type": "regex", "pattern": r"\b(confidential|classified|secret|restricted)\b",
     "tags": ["english", "sensitive"]},
    {"name": "Secreto", "classification": "secreto", "severity": "critical",
     "match_type": "keyword", "pattern": "secreto",
     "tags": ["estado", "secreto"]},
    {"name": "Seguridad del Estado", "classification": "seguridad_del_estado", "severity": "critical",
     "match_type": "keyword", "pattern": "seguridad del estado",
     "tags": ["estado", "seguridad"]},
    {"name": "Restringido", "classification": "restringido", "severity": "high",
     "match_type": "keyword", "pattern": "restringido",
     "tags": ["restringido"]},
]


def _utc_now():
    return datetime.now(timezone.utc).isoformat()


def _state_dir(config_path):
    base = Path(config_path).resolve().parent
    state_dir = base / ".agent_state"
    state_dir.mkdir(parents=True, exist_ok=True)
    try:
        state_dir.chmod(0o700)
    except Exception:
        pass
    return state_dir


def _load_state(path):
    if not path.exists():
        return {}
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return {}


def _save_state(path, data):
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


def _expand_scan_paths(paths):
    expanded = []
    seen = set()
    for path in paths or []:
        value = os.path.expandvars(os.path.expanduser(path))
        if not value:
            continue
        resolved = str(Path(value))
        key = resolved.lower()
        if key in seen:
            continue
        seen.add(key)
        expanded.append(Path(resolved))
    return expanded


def _windows_fixed_drives():
    if os.name != "nt":
        return []
    try:
        import ctypes
        import string
        drive_mask = ctypes.windll.kernel32.GetLogicalDrives()
    except Exception:
        return []
    drives = []
    for index, letter in enumerate(string.ascii_uppercase):
        if not (drive_mask & (1 << index)):
            continue
        root = f"{letter}:\\"
        try:
            drive_type = ctypes.windll.kernel32.GetDriveTypeW(root)
        except Exception:
            continue
        if drive_type in (2, 3):
            drives.append(Path(root))
    return drives


def _windows_scan_roots():
    if os.name != "nt":
        return []
    roots = _windows_fixed_drives()
    return _expand_scan_paths([str(path) for path in roots])


def _linux_scan_roots():
    if os.name == "nt":
        return []
    roots = [Path("/")]
    try:
        result = subprocess.run(
            ["findmnt", "-rno", "TARGET,SOURCE,FSTYPE"],
            capture_output=True, text=True, timeout=10, check=False,
        )
        for line in result.stdout.strip().splitlines():
            parts = line.split()
            if len(parts) >= 2 and parts[1].startswith("/dev/"):
                mount = parts[0]
                if mount not in ("/", "/proc", "/sys", "/dev", "/run", "/boot/efi"):
                    roots.append(Path(mount))
    except Exception:
        pass
    return _expand_scan_paths([str(path) for path in roots])


def _default_paths():
    if os.name == "nt":
        paths = _windows_scan_roots()
    else:
        paths = _linux_scan_roots()
    if not paths:
        paths = _expand_scan_paths([str(Path.home()), "/tmp"])
    return paths


def _default_policy():
    return {
        "code": "aegis-local-shield",
        "name": "Aegis Local Shield",
        "severity": "critical",
        "scan_paths": [str(path) for path in _default_paths()],
        "monitored_extensions": sorted(DEFAULT_DLP_EXTENSIONS),
        "rules": DEFAULT_DLP_KEYWORDS,
    }


def _iter_files(paths, extensions, max_file_size, max_files_per_scan=12000):
    is_win = os.name == "nt"
    excluded_tokens = (
        "\\$recycle.bin",
        "\\system volume information",
        "\\windows\\winsxs",
        "\\windows\\servicing",
        "\\windows\\softwaredistribution\\download",
        "\\programdata\\package cache",
        "\\programdata\\microsoft\\windows\\wer",
        "\\programdata\\microsoft\\windows defender",
    )
    linux_excluded_tokens = (
        "/proc",
        "/sys",
        "/dev",
        "/run",
        "/lost+found",
        "/.snapshots",
    )
    seen = set()
    yielded = 0
    for base in paths:
        try:
            base_exists = base.exists()
        except Exception:
            continue
        if not base_exists:
            continue
        for root, dirs, files in os.walk(base, topdown=True, onerror=lambda exc: logging.getLogger("vant-agent").warning("walk error: %s", exc)):
            root_path = Path(root)
            root_low = str(root_path).lower()
            if not is_win and any(token in root_low for token in linux_excluded_tokens):
                continue
            if is_win and any(token in root_low for token in excluded_tokens):
                continue
            sep = "\\" if is_win else "/"
            if is_win:
                dirs[:] = [
                    d for d in dirs
                    if not any(token in f"{root_low}{sep}{d.lower()}" for token in excluded_tokens)
                ]
            else:
                dirs[:] = [
                    d for d in dirs
                    if not any(token in f"{root_low}{sep}{d.lower()}" for token in linux_excluded_tokens)
                ]
            for name in files:
                path = Path(root) / name
                if str(path) in seen:
                    continue
                seen.add(str(path))
                if extensions and path.suffix.lower() not in extensions:
                    continue
                try:
                    if path.stat().st_size > max_file_size:
                        continue
                except Exception:
                    continue
                yield path
                yielded += 1
                if yielded >= max_files_per_scan:
                    return


def _read_docx_text(path):
    try:
        from docx import Document
        doc = Document(str(path))
        parts = []
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text)
        for section in doc.sections:
            if section.header:
                for p in section.header.paragraphs:
                    if p.text.strip():
                        parts.append(p.text)
            if section.footer:
                for p in section.footer.paragraphs:
                    if p.text.strip():
                        parts.append(p.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        parts.append(cell.text)
        try:
            for comment in doc.comments:
                if comment.text.strip():
                    parts.append(comment.text)
        except Exception:
            pass
        props = doc.core_properties
        for attr in ("title", "subject", "keywords", "category", "comments"):
            val = getattr(props, attr, None)
            if val:
                parts.append(str(val))
        return "\n".join(parts)
    except Exception:
        return ""


def _read_xlsx_text(path):
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), data_only=True, read_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(sheet_name)
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        parts.append(str(cell))
        wb.close()
        return "\n".join(parts)
    except Exception:
        return ""


def _read_pptx_text(path):
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        if p.text.strip():
                            parts.append(p.text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                parts.append(cell.text)
        try:
            for layout in prs.slide_layouts:
                for ph in layout.placeholders:
                    if ph.text.strip():
                        parts.append(ph.text)
        except Exception:
            pass
        return "\n".join(parts)
    except Exception:
        return ""


def _read_odf_text(path):
    suffix = path.suffix.lower()
    try:
        if suffix == ".odt":
            from odf.opendocument import load
            from odf.text import P
            doc = load(str(path))
            parts = []
            for el in doc.getElementsByType(P):
                txt = "".join(node.data for node in el.childNodes if hasattr(node, "data"))
                if txt.strip():
                    parts.append(txt)
            return "\n".join(parts)
        elif suffix == ".ods":
            from odf.opendocument import load
            from odf.table import Table, TableCell
            from odf.text import P as TextP
            doc = load(str(path))
            parts = []
            for table in doc.getElementsByType(Table):
                for cell in table.getElementsByType(TableCell):
                    for p in cell.getElementsByType(TextP):
                        txt = "".join(node.data for node in p.childNodes if hasattr(node, "data"))
                        if txt.strip():
                            parts.append(txt)
            return "\n".join(parts)
        elif suffix == ".odp":
            from odf.opendocument import load
            from odf.draw import Page, Frame
            from odf.text import P as TextP
            doc = load(str(path))
            parts = []
            for page in doc.getElementsByType(Page):
                for frame in page.getElementsByType(Frame):
                    for p in frame.getElementsByType(TextP):
                        txt = "".join(node.data for node in p.childNodes if hasattr(node, "data"))
                        if txt.strip():
                            parts.append(txt)
            return "\n".join(parts)
    except Exception:
        return ""
    return ""


def _read_ole_text(path):
    try:
        import olefile
        if not olefile.isOleFile(str(path)):
            return ""
        ole = olefile.OleFileIO(str(path))
        parts = []
        for stream in ole.listdir():
            stream_path = "/".join(stream)
            try:
                data = ole.openstream(stream_path).read()
                text = data.decode("utf-8", errors="ignore")
                if text.strip():
                    parts.append(text)
            except Exception:
                pass
        ole.close()
        return "\n".join(parts) if parts else ""
    except Exception:
        return ""


def _read_pdf_text(path, max_pages=None):
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        parts = []
        pages_to_process = reader.pages if max_pages is None else reader.pages[:max_pages]
        for page in pages_to_process:
            try:
                text = page.extract_text()
                if text and text.strip():
                    parts.append(text.strip())
            except Exception:
                continue
        return "\n".join(parts)
    except Exception:
        return ""


def _read_file_text(path):
    suffix = path.suffix.lower()
    try:
        if suffix in {".docx", ".docm"}:
            text = _read_docx_text(path)
            if text.strip():
                return text
        elif suffix in {".xlsx", ".xlsm"}:
            text = _read_xlsx_text(path)
            if text.strip():
                return text
        elif suffix in {".pptx", ".pptm"}:
            text = _read_pptx_text(path)
            if text.strip():
                return text
        elif suffix in {".odt", ".ods", ".odp"}:
            text = _read_odf_text(path)
            if text.strip():
                return text
        elif suffix in {".doc", ".xls", ".ppt"}:
            text = _read_ole_text(path)
            if text.strip():
                return text
        elif suffix == ".pdf":
            text = _read_pdf_text(path)
            if text.strip():
                return text
        elif suffix == ".rtf":
            text = path.read_text(encoding="utf-8", errors="ignore")
            if text.strip():
                return text
        elif suffix in {".eml", ".msg"}:
            text = _read_email_text(path)
            if text.strip():
                return text
        else:
            text = path.read_text(encoding="utf-8", errors="ignore")
            if text.strip():
                return text
    except Exception:
        pass
    try:
        data = path.read_bytes()
        text = data.decode("utf-8", errors="ignore")
        chunks = [
            chunk.decode("latin1", errors="ignore")
            for chunk in re.findall(rb"[ -~]{6,}", data)
        ]
        if chunks:
            text = "\n".join([text, *chunks[:800]])
        return text
    except Exception:
        return ""


def _read_email_text(path):
    suffix = path.suffix.lower()
    try:
        if suffix == ".eml":
            import email
            with open(path, "rb") as f:
                msg = email.message_from_binary_file(f)
            parts = []
            if msg["Subject"]:
                parts.append(f"Subject: {msg['Subject']}")
            if msg["From"]:
                parts.append(f"From: {msg['From']}")
            if msg["To"]:
                parts.append(f"To: {msg['To']}")
            if msg["Date"]:
                parts.append(f"Date: {msg['Date']}")
            for part in msg.walk():
                if part.get_content_type().startswith("text/"):
                    try:
                        payload = part.get_payload(decode=True)
                        if payload:
                            parts.append(payload.decode("utf-8", errors="ignore"))
                    except Exception:
                        pass
            return "\n".join(parts)
    except Exception:
        pass
    return ""


def _sha256(path):
    try:
        h = hashlib.sha256()
        with path.open("rb") as handle:
            for chunk in iter(lambda: handle.read(65536), b""):
                h.update(chunk)
        return h.hexdigest()
    except Exception:
        return ""


def _file_owner(path):
    if os.name == "nt":
        try:
            proc = run_hidden(
                ["powershell", "-NoProfile", "-Command",
                 "& {param($p) (Get-Acl -LiteralPath $p).Owner}",
                 "-p", str(path)],
                capture_output=True, text=True, timeout=8, check=False,
            )
            return proc.stdout.strip() if proc.returncode == 0 else ""
        except Exception:
            return ""
    try:
        import pwd
        stat_info = path.stat()
        return pwd.getpwuid(stat_info.st_uid).pw_name
    except Exception:
        return ""


def _path_channel(path):
    low = str(path).lower()
    if "\\downloads\\" in low or "/downloads/" in low:
        return "downloads"
    if "\\desktop\\" in low or "/desktop/" in low:
        return "desktop"
    if "\\documents\\" in low or "/documents/" in low:
        return "documents"
    return "filesystem"


def _metadata_haystack(path, content):
    stat = path.stat() if path.exists() else None
    values = [
        path.name,
        str(path),
        path.suffix.lower(),
        _file_owner(path),
        content,
    ]
    if stat is not None:
        values.extend([
            str(stat.st_size),
            datetime.fromtimestamp(stat.st_ctime, tz=timezone.utc).isoformat(),
            datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc).isoformat(),
        ])
    return "\n".join([item for item in values if item]).lower()


def _detect_new_drives(state, state_path=None):
    if os.name == "nt":
        known = set(state.get("known_drives") or [])
        current = {str(d).lower() for d in _windows_fixed_drives()}
        new_drives = current - known
        state["known_drives"] = sorted(current)
        if new_drives and state_path:
            _save_state(state_path, state)
        return [Path(d) for d in new_drives], state
    known = set(state.get("known_mounts") or [])
    current = set()
    try:
        result = subprocess.run(["mount", "-l"], capture_output=True, text=True, timeout=10, check=False)
        for line in result.stdout.splitlines():
            parts = line.split()
            if len(parts) >= 3 and parts[0].startswith("/dev/"):
                current.add(parts[2].lower())
    except Exception:
        pass
    new_mounts = current - known
    state["known_mounts"] = sorted(current)
    if new_mounts and state_path:
        _save_state(state_path, state)
    return [Path(d) for d in new_mounts], state


class InotifyWatcher:
    """Real-time filesystem watcher using Linux inotify (no pip dependencies)."""

    IN_MODIFY = 0x00000002
    IN_CREATE = 0x00000100
    IN_MOVED_TO = 0x00000080
    IN_DELETE = 0x00000400
    IN_Q_OVERFLOW = 0x00004000
    IN_IS_DIR = 0x40000000
    EVENTS_MASK = IN_MODIFY | IN_CREATE | IN_MOVED_TO | IN_DELETE | IN_Q_OVERFLOW

    def __init__(self, scan_paths, extensions, on_event_callback=None):
        self.logger = logging.getLogger("vant-agent.dlp.watchdog")
        self.scan_paths = scan_paths
        self.extensions = extensions
        self.on_event = on_event_callback
        self._fd = -1
        self._wd_map = {}
        self._running = False
        self._thread = None
        self._new_files_queue = []
        self._lock = threading.Lock()

    def _init_inotify(self):
        try:
            import ctypes
            import ctypes.util
            libc_path = ctypes.util.find_library("c")
            if not libc_path:
                return False
            libc = ctypes.CDLL(libc_path, use_errno=True)
            self._inotify_init = libc.inotify_init
            self._inotify_init.restype = ctypes.c_int
            self._inotify_add_watch = libc.inotify_add_watch
            self._inotify_add_watch.restype = ctypes.c_int
            self._inotify_add_watch.argtypes = [ctypes.c_int, ctypes.c_char_p, ctypes.c_uint32]
            self._inotify_rm_watch = libc.inotify_rm_watch
            self._inotify_rm_watch.restype = ctypes.c_int
            self._inotify_rm_watch.argtypes = [ctypes.c_int, ctypes.c_int]
            self._fd = self._inotify_init()
            if self._fd < 0:
                return False
            return True
        except Exception as exc:
            self.logger.warning("inotify init failed: %s", exc)
            return False

    def _add_watch_recursive(self, path, max_depth=4, current_depth=0):
        if current_depth > max_depth:
            return
        try:
            p = Path(path)
            if not p.exists() or not p.is_dir():
                return
            wd = self._inotify_add_watch(
                self._fd, str(path).encode("utf-8"),
                self.EVENTS_MASK,
            )
            if wd >= 0:
                self._wd_map[wd] = str(path)
        except Exception:
            return
        if current_depth < max_depth:
            try:
                for child in Path(path).iterdir():
                    if child.is_dir() and not child.name.startswith("."):
                        self._add_watch_recursive(child, max_depth, current_depth + 1)
            except Exception:
                pass

    def start(self):
        if os.name == "nt":
            self.logger.info("inotify not available on Windows, using drive-polling fallback")
            return False
        if not self._init_inotify():
            return False
        for path in self.scan_paths:
            self._add_watch_recursive(path, max_depth=3)
        self._running = True
        self._thread = threading.Thread(target=self._watch_loop, daemon=True, name="dlp-inotify")
        self._thread.start()
        self.logger.info("inotify watcher started on %d directories", len(self._wd_map))
        return True

    def stop(self):
        self._running = False
        if self._fd >= 0:
            for wd in self._wd_map:
                try:
                    self._inotify_rm_watch(self._fd, wd)
                except Exception:
                    pass
            try:
                os.close(self._fd)
            except Exception:
                pass
            self._fd = -1
        self._wd_map.clear()

    def _watch_loop(self):
        import select
        while self._running:
            try:
                ready, _, _ = select.select([self._fd], [], [], 2.0)
                if not ready:
                    continue
                data = os.read(self._fd, 65536)
                offset = 0
                while offset < len(data):
                    wd, mask, cookie, name_len = struct.unpack("iIII", data[offset:offset + 16])
                    offset += 16
                    name = data[offset:offset + name_len].rstrip(b"\x00").decode("utf-8", errors="ignore")
                    offset += name_len
                    if mask & self.IN_Q_OVERFLOW:
                        continue
                    if mask & self.IN_IS_DIR:
                        continue
                    if not name:
                        continue
                    ext = Path(name).suffix.lower()
                    if self.extensions and ext not in self.extensions:
                        continue
                    dir_path = self._wd_map.get(wd, "")
                    if dir_path:
                        full_path = os.path.join(dir_path, name)
                        with self._lock:
                            self._new_files_queue.append(full_path)
                        if self.on_event:
                            try:
                                self.on_event(full_path, mask)
                            except Exception:
                                pass
            except Exception as exc:
                if self._running:
                    self.logger.warning("inotify read error: %s", exc)
                time.sleep(1)

    def get_new_files(self):
        with self._lock:
            files = list(self._new_files_queue)
            self._new_files_queue.clear()
        return files


class WindowsFileWatcher:
    """Fallback file watcher for Windows using polling + drive detection."""

    def __init__(self, scan_paths, extensions, on_event_callback=None):
        self.logger = logging.getLogger("vant-agent.dlp.watcher")
        self.scan_paths = scan_paths
        self.extensions = extensions
        self.on_event = on_event_callback
        self._running = False
        self._thread = None
        self._new_files_queue = []
        self._lock = threading.Lock()
        self._snapshot = {}

    def _take_snapshot(self):
        snapshot = {}
        for base in self.scan_paths:
            try:
                for root, dirs, files in os.walk(base, topdown=True, onerror=lambda e: None):
                    for name in files:
                        ext = Path(name).suffix.lower()
                        if self.extensions and ext not in self.extensions:
                            continue
                        fp = os.path.join(root, name)
                        try:
                            st = os.stat(fp)
                            snapshot[fp] = (int(st.st_mtime), st.st_size)
                        except Exception:
                            pass
            except Exception:
                pass
        return snapshot

    def start(self):
        self._snapshot = self._take_snapshot()
        self._running = True
        self._thread = threading.Thread(target=self._watch_loop, daemon=True, name="dlp-winwatch")
        self._thread.start()
        self.logger.info("Windows file watcher started (polling every 10s)")
        return True

    def stop(self):
        self._running = False

    def _watch_loop(self):
        while self._running:
            time.sleep(10)
            try:
                new_snapshot = self._take_snapshot()
                for fp, (mtime, size) in new_snapshot.items():
                    if fp not in self._snapshot or self._snapshot[fp] != (mtime, size):
                        with self._lock:
                            self._new_files_queue.append(fp)
                        if self.on_event:
                            try:
                                self.on_event(fp, 0x00000100)
                            except Exception:
                                pass
                self._snapshot = new_snapshot
            except Exception as exc:
                if self._running:
                    self.logger.warning("winwatch error: %s", exc)

    def get_new_files(self):
        with self._lock:
            files = list(self._new_files_queue)
            self._new_files_queue.clear()
        return files


class AegisDlpService:
    module_name = "aegis_dlp"

    def __init__(self, config_path, cfg):
        self.config_path = config_path
        self.cfg = cfg or {}
        self.state_path = _state_dir(config_path) / "aegis_dlp_state.json"
        self.state = _load_state(self.state_path)
        self.remote_config = {}
        self._init_drive_state()
        self._watcher = None

    def _init_drive_state(self):
        if os.name == "nt":
            known = set(self.state.get("known_drives") or [])
            if not known:
                current = {str(d).lower() for d in _windows_fixed_drives()}
                self.state["known_drives"] = sorted(current)
        else:
            known = set(self.state.get("known_mounts") or [])
            if not known:
                current = set()
                try:
                    result = subprocess.run(["mount", "-l"], capture_output=True, text=True, timeout=10, check=False)
                    for line in result.stdout.splitlines():
                        parts = line.split()
                        if len(parts) >= 3 and parts[0].startswith("/dev/"):
                            current.add(parts[2].lower())
                except Exception:
                    pass
                self.state["known_mounts"] = sorted(current)

    def detect_new_drives(self):
        new, self.state = _detect_new_drives(self.state, self.state_path)
        return new

    def start_watcher(self, scan_paths, extensions):
        if self._watcher is not None:
            return
        def _on_file_event(filepath, mask):
            self.logger.debug("file event: %s mask=0x%x", filepath, mask)
        if os.name == "nt":
            self._watcher = WindowsFileWatcher(scan_paths, extensions, _on_file_event)
        else:
            self._watcher = InotifyWatcher(scan_paths, extensions, _on_file_event)
        self._watcher.start()

    def stop_watcher(self):
        if self._watcher:
            self._watcher.stop()
            self._watcher = None

    def get_realtime_events(self):
        if not self._watcher:
            return []
        return self._watcher.get_new_files()

    def fetch_remote_config(self, control_server, token, agent_id):
        if not control_server:
            return self.remote_config
        try:
            url = f"{control_server}/aegis/api/agent/dlp/config/"
            response = requests.get(
                url,
                headers={"Authorization": f"Bearer {token}"} if token else {},
                timeout=10,
            )
            if response.status_code == 200:
                data = response.json()
                self.remote_config = {"policies": data.get("policies") or []}
        except Exception:
            return self.remote_config
        return self.remote_config

    def submit_threats(self, server_url, token, agent_id, incidents):
        if not server_url or not incidents:
            return None
        max_upload = int(self.cfg.get(self.module_name, {}).get("max_upload_size_mb", 10) or 10)
        max_upload_bytes = max_upload * 1024 * 1024
        enriched = []
        for inc in incidents:
            item = dict(inc)
            file_path = item.get("file_path", "")
            if file_path and max_upload_bytes > 0:
                try:
                    p = Path(file_path)
                    if p.exists() and p.stat().st_size <= max_upload_bytes:
                        item["file_content_base64"] = base64.b64encode(p.read_bytes()).decode("ascii")
                except Exception:
                    pass
            enriched.append(item)
        try:
            response = requests.post(
                f"{server_url}/aegis/api/agent/dlp/threats/",
                json={"agent_id": agent_id, "incidents": enriched},
                headers={"Authorization": f"Bearer {token}"} if token else {},
                timeout=120,
            )
            return response
        except Exception:
            return None

    def queue_incidents(self, incidents):
        pending = self.state.get("pending_incidents") or []
        existing = {item.get("fingerprint", "") for item in pending}
        for incident in incidents or []:
            fingerprint = incident.get("fingerprint", "")
            if not fingerprint or fingerprint in existing:
                continue
            pending.append(incident)
            existing.add(fingerprint)
        self.state["pending_incidents"] = pending[-5000:]
        _save_state(self.state_path, self.state)
        return self.state["pending_incidents"]

    def peek_pending_incidents(self):
        return list(self.state.get("pending_incidents") or [])

    def clear_pending_incidents(self):
        self.state["pending_incidents"] = []
        _save_state(self.state_path, self.state)

    def scan(self):
        dlp_cfg = self.cfg.get(self.module_name, {}) or {}
        if not dlp_cfg.get("enabled", True):
            return []

        policies = (self.remote_config or {}).get("policies") or []
        if not policies:
            policies = [_default_policy()]

        rules = []
        scan_paths = []
        default_paths = [str(path) for path in _default_paths()]
        monitored_extensions = set()
        max_file_size_mb = int(dlp_cfg.get("max_file_size_mb", 25) or 25)
        max_files_per_scan = int(dlp_cfg.get("max_files_per_scan", 12000) or 12000)
        max_scan_seconds = int(dlp_cfg.get("max_scan_seconds", 0) or 0)
        realtime_enabled = dlp_cfg.get("realtime_enabled", True)

        for policy in policies:
            scan_mode = policy.get("scan_mode", "all")
            if scan_mode == "all":
                scan_paths = [str(path) for path in _default_paths()]
            else:
                scan_paths.extend(policy.get("scan_paths") or [])
            monitored_extensions.update([ext.lower() for ext in (policy.get("monitored_extensions") or [])])
            max_file_size_mb = max(max_file_size_mb, int(policy.get("max_file_size_mb", 10) or 10))
            for rule in policy.get("rules") or []:
                rules.append((policy, rule))

        new_drives = self.detect_new_drives()
        if new_drives:
            for d in new_drives:
                scan_paths.append(str(d))
                self.logger.info("new drive detected: %s, adding to scan", d)

        local_scan_paths = dlp_cfg.get("scan_paths") or []
        if local_scan_paths:
            scan_paths = list(local_scan_paths) + scan_paths

        if not scan_paths:
            scan_paths = default_paths
        else:
            for path in default_paths:
                if path not in scan_paths:
                    scan_paths.append(path)

        monitored_extensions.update([ext.lower() for ext in (dlp_cfg.get("monitored_extensions") or [])])
        monitored_extensions.update(DEFAULT_DLP_EXTENSIONS)
        paths = _expand_scan_paths(scan_paths)
        max_file_size = max_file_size_mb * 1024 * 1024

        if realtime_enabled and self._watcher is None:
            self.start_watcher(scan_paths, monitored_extensions)

        realtime_files = self.get_realtime_events()

        incidents = []
        known = set(self.state.get("incident_keys", []))
        scan_cache = self.state.get("scan_cache") or {}
        updated_cache = {}
        started_at = time.monotonic()

        files_to_scan = []
        for path in _iter_files(
            paths, monitored_extensions, max_file_size, max_files_per_scan=max_files_per_scan
        ):
            if max_scan_seconds > 0 and (time.monotonic() - started_at) >= max_scan_seconds:
                break
            files_to_scan.append(path)

        for rf in realtime_files:
            rp = Path(rf)
            if rp.exists() and rp.suffix.lower() in monitored_extensions:
                try:
                    if rp.stat().st_size <= max_file_size:
                        files_to_scan.append(rp)
                except Exception:
                    pass

        for path in files_to_scan:
            try:
                stat = path.stat()
                cache_value = f"{int(stat.st_mtime)}:{stat.st_size}"
            except Exception:
                continue
            cache_key = str(path).lower()
            updated_cache[cache_key] = cache_value
            if scan_cache.get(cache_key) == cache_value:
                continue
            content = _read_file_text(path)
            haystack = _metadata_haystack(path, content)
            if not haystack.strip():
                continue
            for policy, rule in rules:
                pattern = (rule.get("pattern") or "").strip()
                if not pattern:
                    continue
                matched_terms = []
                match_type = (rule.get("match_type") or "keyword").strip().lower()
                if match_type == "regex":
                    if re.search(pattern, haystack, flags=re.IGNORECASE):
                        matched_terms.append(pattern)
                elif match_type == "metadata":
                    if pattern.lower() in haystack:
                        matched_terms.append(pattern)
                else:
                    if pattern.lower() in haystack:
                        matched_terms.append(pattern)
                if not matched_terms:
                    continue

                file_hash = _sha256(path)
                incident_key = hashlib.sha256(
                    f"{file_hash}|{path}|{policy.get('code')}|{rule.get('name')}".encode("utf-8")
                ).hexdigest()
                if incident_key in known:
                    continue
                known.add(incident_key)
                incidents.append({
                    "fingerprint": incident_key,
                    "policy_code": policy.get("code", ""),
                    "rule_name": rule.get("name", ""),
                    "classification": rule.get("classification") or policy.get("code", ""),
                    "severity": rule.get("severity") or policy.get("severity", "high"),
                    "file_name": path.name,
                    "file_path": str(path),
                    "file_hash": file_hash,
                    "actor": _file_owner(path),
                    "channel": _path_channel(path),
                    "status": "open",
                    "detected_at": _utc_now(),
                    "summary": f"Coincidencia DLP en {path.name}",
                    "matched_keywords": matched_terms,
                    "metadata": {
                        "size": path.stat().st_size if path.exists() else 0,
                        "suffix": path.suffix.lower(),
                        "created_at": datetime.fromtimestamp(
                            path.stat().st_ctime, tz=timezone.utc
                        ).isoformat() if path.exists() else "",
                        "modified_at": datetime.fromtimestamp(
                            path.stat().st_mtime, tz=timezone.utc
                        ).isoformat() if path.exists() else "",
                    },
                })
                if len(incidents) >= 25:
                    self.queue_incidents(incidents)
                    incidents = []
        self.queue_incidents(incidents)
        self.state["incident_keys"] = sorted(list(known))[-5000:]
        self.state["scan_cache"] = dict(list(updated_cache.items())[-25000:])
        _save_state(self.state_path, self.state)
        return incidents
